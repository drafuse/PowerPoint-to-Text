import argparse
from pathlib import Path
import tkinter as tk
from tkinter import filedialog, messagebox
import webbrowser

from pptx import Presentation


def extract_pptx_text(pptx_path: Path, output_txt: Path) -> None:
    prs = Presentation(pptx_path)
    lines = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_number} ---")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append(text)

        lines.append("")

    output_txt.write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract text from a .pptx file.")
    parser.add_argument("pptx_path", type=Path, nargs="?", help="Path to the input .pptx file")
    parser.add_argument(
        "-o",
        "--output",
        dest="output_txt",
        type=Path,
        help="Path to output .txt file (default: same name as input)",
    )
    parser.add_argument("--gui", action="store_true", help="Launch window interface")
    return parser.parse_args()


def validate_input_path(pptx_path: Path) -> str | None:
    if not pptx_path.exists():
        return f"Input file not found: {pptx_path}"
    if pptx_path.suffix.lower() != ".pptx":
        return f"Input must be a .pptx file: {pptx_path}"
    return None


def run_cli(pptx_path: Path, output_txt: Path | None) -> int:
    input_path = pptx_path.expanduser().resolve()
    error = validate_input_path(input_path)
    if error:
        print(error)
        return 1

    output_path = (output_txt or input_path.with_suffix(".txt")).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    extract_pptx_text(input_path, output_path)
    print(f"Text exported successfully to: {output_path}")
    return 0


def run_gui() -> int:
    root = tk.Tk()
    root.title("PowerPoint to Text")
    root.geometry("640x280")
    root.resizable(False, False)

    input_var = tk.StringVar()
    output_var = tk.StringVar()

    def browse_input() -> None:
        selected = filedialog.askopenfilename(
            title="Select PowerPoint file",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
        )
        if selected:
            input_var.set(selected)
            if not output_var.get().strip():
                output_var.set(str(Path(selected).with_suffix(".txt")))

    def browse_output() -> None:
        selected = filedialog.asksaveasfilename(
            title="Select output text file",
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt"), ("All files", "*.*")],
        )
        if selected:
            output_var.set(selected)

    def convert() -> None:
        input_text = input_var.get().strip()
        output_text = output_var.get().strip()

        if not input_text:
            messagebox.showerror("Missing Input", "Please select a .pptx file.")
            return

        input_path = Path(input_text).expanduser().resolve()
        error = validate_input_path(input_path)
        if error:
            messagebox.showerror("Invalid Input", error)
            return

        output_path = (
            Path(output_text).expanduser().resolve()
            if output_text
            else input_path.with_suffix(".txt")
        )

        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            extract_pptx_text(input_path, output_path)
            messagebox.showinfo("Success", f"Text exported successfully to:\n{output_path}")
        except Exception as exc:
            messagebox.showerror("Conversion Failed", str(exc))

    pad = {"padx": 12, "pady": 8}

    tk.Label(root, text="PowerPoint (.pptx):").grid(row=0, column=0, sticky="w", **pad)
    tk.Entry(root, textvariable=input_var, width=62).grid(row=1, column=0, sticky="we", **pad)
    tk.Button(root, text="Browse...", width=12, command=browse_input).grid(row=1, column=1, **pad)

    tk.Label(root, text="Output Text (.txt):").grid(row=2, column=0, sticky="w", **pad)
    tk.Entry(root, textvariable=output_var, width=62).grid(row=3, column=0, sticky="we", **pad)
    tk.Button(root, text="Save As...", width=12, command=browse_output).grid(row=3, column=1, **pad)

    tk.Button(root, text="Convert", width=16, command=convert).grid(row=4, column=0, pady=16)

    credits = tk.Frame(root)
    credits.grid(row=5, column=0, sticky="w", padx=12, pady=(0, 10))

    tk.Label(credits, text="Created by Dan Rafuse - ").pack(side="left")

    github = tk.Label(credits, text="github.com/drafuse", fg="blue", cursor="hand2")
    github.pack(side="left")
    github.bind("<Button-1>", lambda _e: webbrowser.open("https://github.com/drafuse"))

    tk.Label(credits, text=" using ").pack(side="left")

    pycharm = tk.Label(credits, text="PyCharm", fg="blue", cursor="hand2")
    pycharm.pack(side="left")
    pycharm.bind("<Button-1>", lambda _e: webbrowser.open("https://www.jetbrains.com/pycharm/"))

    tk.Label(credits, text=" and ").pack(side="left")

    codex = tk.Label(credits, text="Codex", fg="blue", cursor="hand2")
    codex.pack(side="left")
    codex.bind("<Button-1>", lambda _e: webbrowser.open("https://openai.com/codex/"))

    root.mainloop()
    return 0


def main() -> int:
    args = parse_args()
    if args.gui or args.pptx_path is None:
        return run_gui()
    return run_cli(args.pptx_path, args.output_txt)


if __name__ == "__main__":
    raise SystemExit(main())
